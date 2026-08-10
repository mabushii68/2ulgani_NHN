# -*- coding: utf-8 -*-
# 러다이트 2026 발표 PPT 생성기 — 14장, 16:9, 게임 색 계약(마젠타 = AI 위협 전용) 준수
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

BG      = RGBColor(0x0A, 0x0A, 0x12)   # 거의 검정
PANEL   = RGBColor(0x14, 0x14, 0x20)   # 패널
TEXT    = RGBColor(0xEB, 0xEB, 0xF0)   # 본문
DIM     = RGBColor(0x8C, 0x92, 0x9E)   # 보조
MAGENTA = RGBColor(0xFF, 0x2E, 0xB8)   # AI 위협 전용 강조
BLUE    = RGBColor(0x4D, 0xA6, 0xFF)   # 문과
GREEN   = RGBColor(0x53, 0xC8, 0x7A)   # 이과
YELLOW  = RGBColor(0xE8, 0xC5, 0x4A)   # 예체능
FONT    = "맑은 고딕"

W, H = Inches(13.333), Inches(7.5)
prs = Presentation()
prs.slide_width, prs.slide_height = W, H
BLANK = prs.slide_layouts[6]

def slide():
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = BG
    return s

def box(s, x, y, w, h, fill=None):
    from pptx.enum.shapes import MSO_SHAPE
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill or PANEL
    sh.line.fill.background()
    sh.shadow.inherit = False
    return sh

def text(s, x, y, w, h, runs, size=18, color=TEXT, bold=False, align=PP_ALIGN.LEFT,
         anchor=MSO_ANCHOR.TOP, line_spacing=1.15):
    """runs: str 또는 [(문구, 색, 볼드)] 리스트의 리스트(문단들)"""
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    if isinstance(runs, str):
        runs = [[(runs, color, bold)]]
    first = True
    for para in runs:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        p.line_spacing = line_spacing
        if isinstance(para, str):
            para = [(para, color, bold)]
        for t, c, b in para:
            r = p.add_run()
            r.text = t
            r.font.name = FONT
            r.font.size = Pt(size)
            r.font.color.rgb = c
            r.font.bold = b
    return tb

def shot(s, x, y, w, h, label):
    """스크린샷 자리 표시 박스"""
    b = box(s, x, y, w, h, fill=RGBColor(0x1B, 0x1B, 0x2A))
    tf = b.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = "[ 스크린샷 ]\n" + label
    r.font.name = FONT; r.font.size = Pt(13); r.font.color.rgb = DIM

def header(s, title, accent=False):
    text(s, 0.7, 0.45, 11.9, 0.9, title, size=30, bold=True,
         color=MAGENTA if accent else TEXT)

# ── 1. 타이틀 ──
s = slide()
text(s, 0.7, 2.1, 11.9, 1.3, "러다이트 2026", size=54, bold=True, align=PP_ALIGN.CENTER)
text(s, 0.7, 3.4, 11.9, 0.7, [[("AI는 당신의 플레이를 학습합니다. ", TEXT, False),
                                ("그러니 AI에게 거짓말하세요.", MAGENTA, True)]],
     size=22, align=PP_ALIGN.CENTER)
text(s, 0.7, 5.6, 11.9, 1.2, [
    "NHN 게임 AI 공모전 · 팀 2인 (이양빈 — 기획/AI/전투 · 김정준 — 아트/UI/파이프라인)",
    "플레이: https://mabushii68.github.io/2ulgani_NHN/  ·  개발 7일 (08.04 – 08.10)"],
    size=14, color=DIM, align=PP_ALIGN.CENTER)

# ── 2. 역발상 ──
s = slide(); header(s, "게임 AI의 역발상")
text(s, 0.7, 1.8, 11.9, 1.0, "대부분의 게임 AI는  숨겨져 있고,  더 강해지는 방향으로 작동한다.", size=20, color=DIM)
text(s, 0.7, 2.9, 11.9, 1.6, [[("우리는 AI를 ", TEXT, False), ("공개", MAGENTA, True),
                                ("하고, 플레이어가 ", TEXT, False), ("속이는 대상", MAGENTA, True),
                                ("으로 만들었다.", TEXT, False)]], size=26, bold=False)
text(s, 0.7, 4.9, 11.9, 1.6, [
    "AI의 학습 현황(확률·신뢰도)이 화면에 항상 떠 있다 — 숨겨진 치트가 아니라 읽을 수 있는 상대.",
    "따라서 이 게임의 실력은 조작 숙련이 아니라, 자신에 대한 예측을 읽고 배신하는 심리전이다."],
    size=16, color=DIM)

# ── 3. 핵심 루프 ──
s = slide(); header(s, "핵심 루프 — 학습 → 공개 → 예측 → 배신", accent=True)
steps = [
    ("① 관찰", "적탄이 접근할 때마다 좌/우\n회피를 기록 (피격 위기 이벤트)", TEXT),
    ("② 공개", "AI MODEL: LEARNING…\n→ LEFT 73% [HIGH]", TEXT),
    ("③ 예측탄", "확신한 AI가 당신이 '피할 자리'로\n마젠타 예측탄 발사", MAGENTA),
    ("④ 배신", "조준선을 읽고 반대로 회피\n→ PREDICTION FAILED", MAGENTA),
]
for i, (t1, t2, c) in enumerate(steps):
    x = 0.7 + i * 3.12
    box(s, x, 1.9, 2.85, 2.2)
    text(s, x + 0.2, 2.1, 2.45, 0.6, t1, size=20, bold=True, color=c)
    text(s, x + 0.2, 2.75, 2.45, 1.3, t2, size=13, color=DIM)
shot(s, 0.7, 4.5, 12.0, 2.4, "②→④ 연속 장면: HUD HIGH 전환 → 마젠타 조준선 → PREDICTION FAILED 연출")

# ── 4. 읽을 수 있는 AI (UI) ──
s = slide(); header(s, "읽을 수 있는 AI — 확률 표기가 곧 공략 정보")
text(s, 0.7, 1.5, 5.6, 4.0, [
    [("HUD 미니 패널", TEXT, True)],
    [("상시 노출: 학습 상태 · 우세 방향 · 신뢰도", DIM, False)],
    [("", DIM, False)],
    [("TARGET PROFILE (상자 오픈 화면)", TEXT, True)],
    [("웨이브마다 AI가 본 나를 열람 — 표본 수, 좌/우 편향,", DIM, False)],
    [("예측 적중률. 다음 방 전에 '나를 어떻게 속일지' 계획", DIM, False)],
    [("", DIM, False)],
    [("COUNTER PROTOCOL", TEXT, True)],
    [("AI를 겨냥한 업그레이드 2종 — 행동교정(관측 삭제),", DIM, False)],
    [("논문조작(가짜 표본 주입으로 역방향 학습 유도)", DIM, False)]], size=16)
shot(s, 6.6, 1.5, 6.1, 2.5, "HUD 미니 패널 (LEFT 73% HIGH 상태)")
shot(s, 6.6, 4.2, 6.1, 2.5, "TARGET PROFILE 패널")

# ── 5. 진행 구조 ──
s = slide(); header(s, "진행 구조 — 선형 던전 + 세부전공")
text(s, 0.7, 1.5, 11.9, 0.9, "시작방 ─복도─ 전투방 1 ~ 6 ─복도─ 보스방   (방 락인: 전멸해야 문이 열린다 — 시간 버티기 불가)",
     size=16, color=DIM)
text(s, 0.7, 2.5, 11.9, 0.6, [[("전공 3 × 세부전공 3 = 9가지 정체성, 9가지 탄막", TEXT, True)]], size=20)
majors = [("문과", BLUE, "어문계=연필 · 상경계=돈 · 법조계=책"),
          ("이과", GREEN, "자연과학=숫자 · 공학=번개 · 컴퓨터과학=모니터"),
          ("예체능", YELLOW, "체육=공 · 미술=붓 · 음악=음표")]
for i, (m, c, d) in enumerate(majors):
    x = 0.7 + i * 4.1
    box(s, x, 3.3, 3.85, 1.5)
    text(s, x + 0.25, 3.5, 3.35, 0.5, m, size=18, bold=True, color=c)
    text(s, x + 0.25, 4.05, 3.35, 0.6, d, size=12, color=DIM)
shot(s, 0.7, 5.1, 12.0, 1.9, "세부전공 선택 화면 + 테마 탄막 발사 장면 (2장 나란히)")

# ── 6. 학습 모델 (기술 1) ──
s = slide(); header(s, "AI 구조 ① — 무엇을 학습하는가", accent=True)
text(s, 0.7, 1.5, 11.9, 1.0, [[("피격 위기 이벤트: ", TEXT, True),
    ("TTI(충돌까지 시간) 0.5초 이내 탄 접근 시 트리거 → 0.6초 뒤 좌/우 회피 판정 (변위 0.3u 미만은 표본 제외)", DIM, False)]], size=17)
text(s, 0.7, 2.7, 11.9, 1.0, [[("LEFT / RIGHT 2분류 조건부 확률 — ", TEXT, True),
    ("의도적 절제: 8방향 회귀가 아니라 2분류여야 플레이어가 '읽고 속일 수 있는' 단위가 된다", DIM, False)]], size=17)
text(s, 0.7, 4.0, 11.9, 1.7, [
    [("수식: ", TEXT, True), ("P(LEFT) = (n_left + 1) / (n_total + 2)   — 가상 카운트(1,1)로 초기 과확신 방지", DIM, False)],
    [("망각: ", TEXT, True), ("웨이브 종료마다 관측 카운트 ×0.8 — 옛 습관은 흐려지고 최근 행동이 지배", DIM, False)],
    [("게이트: ", TEXT, True), ("표본 수 + 편향 강도 임계를 넘어야 신뢰도 HIGH → 그때만 예측탄 발사", DIM, False)]], size=16)
text(s, 0.7, 6.2, 11.9, 0.8, "구현: AIBrain은 순수 C# (엔진 무의존) — 가짜 이벤트 시퀀스로 단위 검증, MonoBehaviour 어댑터가 이벤트만 전달",
     size=14, color=DIM)

# ── 7. 심리전 (기술 2) ──
s = slide(); header(s, "AI 구조 ② — 예측탄과 배신", accent=True)
text(s, 0.7, 1.5, 6.0, 4.6, [
    [("예측탄", MAGENTA, True)],
    [("신뢰도 HIGH일 때 엘리트·보스가 플레이어의", DIM, False)],
    [("'회피 예상 지점'을 조준 — 마젠타 조준선을 먼저", DIM, False)],
    [("보여준다 (읽을 시간을 준다 = 심리전 성립 조건)", DIM, False)],
    [("", DIM, False)],
    [("PREDICTION FAILED", MAGENTA, True)],
    [("예측을 배신하면 히트스톱 + 글리치 연출 +", DIM, False)],
    [("신뢰도 하락이 눈앞에서 벌어진다", DIM, False)],
    [("", DIM, False)],
    [("역카운터", TEXT, True)],
    [("예측탄의 예측 방향 vs 실제 회피를 대조 집계 —", DIM, False)],
    [("결과 프로필의 '역카운터 성공률'로 환산", DIM, False)]], size=16)
shot(s, 7.0, 1.5, 5.7, 2.5, "마젠타 조준선 + 예측탄")
shot(s, 7.0, 4.2, 5.7, 2.5, "PREDICTION FAILED 글리치 연출")

# ── 8. 보스 P2 ──
s = slide(); header(s, "보스 「거대 LLM」 P2 — PATTERN: YOU", accent=True)
text(s, 0.7, 1.6, 11.9, 0.9, [[("게임 내내 나를 학습한 AI가, 마지막에 ", TEXT, False),
                                ("나 자신이 되어 돌아온다", MAGENTA, True)]], size=24)
text(s, 0.7, 2.8, 11.9, 1.9, [
    "· 교전 거리 복제 — 내가 유지하던 거리를 보스가 유지한다",
    "· 무기 복제 — 내 전공 무기의 마젠타 버전을 쏜다",
    "· 구역 선점 — 내가 좋아하던 자리에 장판을 깐다",
    "· 예측탄 — 축적된 학습의 총화"], size=17, color=DIM)
shot(s, 0.7, 4.9, 12.0, 2.1, "보스 P2 진입 오버레이 + 마젠타 오라 전투 장면")

# ── 9. 결과 화면 ──
s = slide(); header(s, "최종 점수는 숫자가 아니라 프로필이다")
text(s, 0.7, 1.6, 5.8, 4.4, [
    [("결과 화면 「AI가 학습한 나」", TEXT, True)],
    [("", DIM, False)],
    [("· AI가 붙여준 별명 (플레이 스타일 기반)", DIM, False)],
    [("· 평균 교전 거리 / 회피 좌우 편향", DIM, False)],
    [("· AI 예측 적중률 vs 역카운터 성공률", DIM, False)],
    [("", DIM, False)],
    [("승패와 무관하게 출력 — 패배해도", TEXT, False)],
    [("'AI에게 어떻게 읽혔는가'가 남는다", TEXT, False)]], size=17)
shot(s, 6.8, 1.6, 5.9, 4.4, "결과 프로필 화면 (별명 + 통계)")

# ── 10. 세부 설계 판단 ──
s = slide(); header(s, "세부 설계 판단들")
text(s, 0.7, 1.6, 11.9, 4.8, [
    [("매크로 DDA — ", TEXT, True), ("잘하는 판엔 스폰 가속, 몰리는 판엔 감속. AI 학습과 분리된 별도 축", DIM, False)],
    [("", DIM, False)],
    [("예측탄 온스크린 게이트 — ", TEXT, True), ("화면 밖에서 쏘는 예측탄은 심리전이 아니라 이지선다 강요 → 화면 안에서만", DIM, False)],
    [("", DIM, False)],
    [("색 계약 — ", TEXT, True), ("마젠타·핫핑크 = 'AI가 나를 읽고 행하는 것' 전용 예약. 그 외 위협은 무채색~주황", DIM, False)],
    [("", DIM, False)],
    [("실루엣 차별 — ", TEXT, True), ("위협 신호를 색 하나에 걸지 않는다 (색맹 접근성, 예측탄은 모양도 다름)", DIM, False)],
    [("", DIM, False)],
    [("전멸형 종료 — ", TEXT, True), ("시간제면 '버티기'가 최적해가 되어 AI에게 줄 회피 표본 자체가 사라진다", DIM, False)]], size=16)

# ── 11. AI로 만든 게임 ──
s = slide(); header(s, "개발 프로세스 — AI로 AI 게임 만들기", accent=True)
text(s, 0.7, 1.6, 11.9, 0.8, [[("Claude Code + Unity MCP", TEXT, True),
                                ("  —  코드 작성 · 에디터 조작 · 씬 구성 · 테스트 · 문서화 전 과정", DIM, False)]], size=19)
stats = [("7일", "개발 기간"), ("2인", "팀"), ("106+", "커밋"), ("46+", "AI 세션 로그"), ("18.2MB", "WebGL 빌드")]
for i, (n, l) in enumerate(stats):
    x = 0.7 + i * 2.5
    box(s, x, 2.7, 2.25, 1.5)
    text(s, x, 2.9, 2.25, 0.7, n, size=26, bold=True, align=PP_ALIGN.CENTER)
    text(s, x, 3.6, 2.25, 0.5, l, size=13, color=DIM, align=PP_ALIGN.CENTER)
text(s, 0.7, 4.6, 11.9, 2.2, [
    "· 모든 세션을 AI_USAGE_LOG에 실시간 기록 — 프롬프트 원문, 결과, 실패까지 (사후 재구성 없음)",
    "· 씬·프리팹·SO 구성은 전부 재실행 가능한 결정론 에디터 빌더로 — 사람 실수를 코드 재실행으로 복구",
    "· 실패 기록 보존: ML-Agents 미착수 판단, 도구 한계 우회 등 — '시도와 판단'의 서사가 기술 문서의 자산",
    "· 탄막 아트 5종은 절차 생성(문자열 마스크 픽셀 아트), SFX·BGM 10종은 수식 신스로 자체 합성"], size=14, color=DIM)

# ── 12. 제원 ──
s = slide(); header(s, "제원")
text(s, 0.7, 1.6, 11.9, 4.6, [
    [("엔진  ", DIM, False), ("Unity 6000.3.7f1 · URP 2D Renderer · 단일 씬 상태 머신", TEXT, False)],
    [("배포  ", DIM, False), ("WebGL 18.24MB · Gzip · GitHub Pages — 링크 클릭만으로 플레이, 설치·로그인 없음", TEXT, False)],
    [("AI    ", DIM, False), ("온라인 조건부 확률 모델 (순수 C#) — 런타임 학습, 사전 훈련 없음, 서버 없음", TEXT, False)],
    [("에셋  ", DIM, False), ("Franuka 픽셀 아트 팩 (라이선스 준수 · CREDITS 전량 기재) + 자체 절차 생성분", TEXT, False)],
    [("문서  ", DIM, False), ("게임 소개 · AI 활용 기술 문서 · 팀 롤 기술서 — 매 세션 로그가 원천", TEXT, False)]], size=17, line_spacing=1.6)

# ── 13. 데모 ──
s = slide()
text(s, 0.7, 3.0, 11.9, 1.0, "플레이 영상", size=40, bold=True, align=PP_ALIGN.CENTER)
text(s, 0.7, 4.2, 11.9, 0.6, "(30~60초 제출 영상 삽입 — 김정준)", size=16, color=DIM, align=PP_ALIGN.CENTER)

# ── 14. 마무리 ──
s = slide()
text(s, 0.7, 2.4, 11.9, 1.0, [[("당신은 AI에게 ", TEXT, False), ("어떻게 기억될", MAGENTA, True),
                                (" 것인가", TEXT, False)]], size=36, bold=False, align=PP_ALIGN.CENTER)
text(s, 0.7, 4.4, 11.9, 1.4, [
    "플레이  https://mabushii68.github.io/2ulgani_NHN/",
    "소스  https://github.com/mabushii68/2ulgani_NHN"], size=16, color=DIM, align=PP_ALIGN.CENTER)

import os
out = r"C:\Users\Vin Jacob Lee\Desktop\dd\2ulgani_NHN\Assets\_Project\Docs\PPT~"
os.makedirs(out, exist_ok=True)
path = os.path.join(out, "러다이트2026_발표.pptx")
prs.save(path)
print("저장:", path, os.path.getsize(path), "bytes,", len(prs.slides.__iter__.__self__._sldIdLst), "slides")
